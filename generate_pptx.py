#!/usr/bin/env python3
"""
Generate Accesly Pitch Deck PPTX
Respects the design elements from the HTML version
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors from the HTML design
COLORS = {
    'bg': RgbColor(0x00, 0x00, 0x00),
    'lavender': RgbColor(0x8B, 0x6C, 0xE7),
    'lavender_light': RgbColor(0xA9, 0x8D, 0xF0),
    'mint': RgbColor(0x45, 0xC9, 0xA8),
    'blush': RgbColor(0xFF, 0x8F, 0xAB),
    'white': RgbColor(0xFF, 0xFF, 0xFF),
    'gray3': RgbColor(0xD4, 0xCE, 0xD9),
    'gray4': RgbColor(0x9E, 0x95, 0xA7),
    'gray5': RgbColor(0x6B, 0x5F, 0x78),
    'dark_card': RgbColor(0x1A, 0x1A, 0x2E),
}

def set_slide_background(slide, color):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, 
                 font_color=None, bold=False, font_name='Arial'):
    """Add a text box with styling"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.name = font_name
    if font_color:
        p.font.color.rgb = font_color
    return txBox

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Blank layout
    blank_layout = prs.slide_layouts[6]
    
    # ===================
    # SLIDE 1: Title
    # ===================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, COLORS['bg'])
    
    # Title
    add_text_box(slide1, Inches(0.5), Inches(2.5), Inches(12), Inches(1.5),
                 "Accesly", font_size=72, font_color=COLORS['lavender'], bold=True)
    
    # Subtitle
    add_text_box(slide1, Inches(0.5), Inches(4), Inches(12), Inches(1),
                 "Focus on your value proposition.\nWe handle the friction.",
                 font_size=32, font_color=COLORS['white'], bold=True)
    
    # ===================
    # SLIDE 2: Problem
    # ===================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2, COLORS['bg'])
    
    # Tag
    add_text_box(slide2, Inches(0.5), Inches(0.5), Inches(3), Inches(0.5),
                 "THE PROBLEM", font_size=12, font_color=COLORS['lavender'], bold=True)
    
    # Main text
    add_text_box(slide2, Inches(0.5), Inches(1.2), Inches(10), Inches(1.5),
                 "Web3 security puts all responsibility on the user and users fail.",
                 font_size=36, font_color=COLORS['white'], bold=True)
    
    # Stats boxes - we'll describe them as text
    stats = [
        ("$1.7B", "Lost to key compromises", "CertiK, H1 2025", COLORS['lavender']),
        ("69%", "From wallet compromises", "CertiK, H1 2025", COLORS['blush']),
        ("3.7M", "BTC lost to seed phrases", "Ledger, 2025", COLORS['mint']),
    ]
    
    for i, (num, desc, source, color) in enumerate(stats):
        left = Inches(0.5 + i * 4)
        # Number
        add_text_box(slide2, left, Inches(3.5), Inches(3.5), Inches(1),
                     num, font_size=48, font_color=color, bold=True)
        # Description
        add_text_box(slide2, left, Inches(4.5), Inches(3.5), Inches(0.5),
                     desc, font_size=14, font_color=COLORS['gray3'])
        # Source
        add_text_box(slide2, left, Inches(5), Inches(3.5), Inches(0.5),
                     source, font_size=10, font_color=COLORS['gray5'])
    
    # ===================
    # SLIDE 3: Solution
    # ===================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3, COLORS['bg'])
    
    add_text_box(slide3, Inches(0.5), Inches(0.5), Inches(3), Inches(0.5),
                 "THE SOLUTION", font_size=12, font_color=COLORS['mint'], bold=True)
    
    add_text_box(slide3, Inches(0.5), Inches(1.2), Inches(10), Inches(1),
                 "One click. 3 seconds. Your wallet is ready.",
                 font_size=36, font_color=COLORS['white'], bold=True)
    
    # Flow description
    flow = """
    1️⃣ SIGN IN
    Sign in with Google or Email
    
    ➡️
    
    2️⃣ WALLET CREATED
    Your Stellar wallet is ready
    Public Key: GAXK...7F2M
    
    ➡️
    
    3️⃣ READY TO TRANSACT
    Balance: 0.00 XLM
    Send | Receive
    """
    add_text_box(slide3, Inches(0.5), Inches(2.5), Inches(12), Inches(4),
                 flow, font_size=18, font_color=COLORS['gray3'])
    
    # Powered by Etherfuse
    add_text_box(slide3, Inches(5), Inches(6.5), Inches(4), Inches(0.5),
                 "Powered by Etherfuse", font_size=14, font_color=COLORS['gray4'])
    
    # ===================
    # SLIDE 4: Live Demo
    # ===================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4, COLORS['bg'])
    
    add_text_box(slide4, Inches(0.5), Inches(0.5), Inches(3), Inches(0.5),
                 "LIVE DEMO", font_size=12, font_color=COLORS['lavender'], bold=True)
    
    add_text_box(slide4, Inches(2), Inches(3), Inches(9), Inches(1.5),
                 "See it in action.", font_size=48, font_color=COLORS['white'], bold=True)
    
    add_text_box(slide4, Inches(3), Inches(5), Inches(7), Inches(0.5),
                 "[Split screen: User view | Developer view]",
                 font_size=18, font_color=COLORS['gray4'])
    
    # ===================
    # SLIDE 5: Traction
    # ===================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5, COLORS['bg'])
    
    add_text_box(slide5, Inches(0.5), Inches(0.5), Inches(3), Inches(0.5),
                 "TRACTION", font_size=12, font_color=COLORS['mint'], bold=True)
    
    add_text_box(slide5, Inches(0.5), Inches(1.2), Inches(6), Inches(1),
                 "Growing on Stellar", font_size=36, font_color=COLORS['white'], bold=True)
    
    add_text_box(slide5, Inches(0.5), Inches(1.8), Inches(6), Inches(0.5),
                 "Key metrics and milestones", font_size=14, font_color=COLORS['gray4'])
    
    # Chart description
    add_text_box(slide5, Inches(0.5), Inches(2.5), Inches(12), Inches(1.5),
                 "📈 Wallets Created - Real-time growth chart from Feb 10 to Mar 22",
                 font_size=18, font_color=COLORS['gray3'])
    
    # Stats
    traction_stats = [
        ("112", "Wallets created", "Real-time from Stellar"),
        ("3s", "Onboarding time", "Average user flow"),
        ("396", "npm Downloads", "First week"),
        ("39", "Active integrations", "Developer projects"),
    ]
    
    for i, (num, title, sub) in enumerate(traction_stats):
        left = Inches(0.5 + i * 3.1)
        add_text_box(slide5, left, Inches(4.5), Inches(2.8), Inches(0.8),
                     num, font_size=36, font_color=COLORS['white'], bold=True)
        add_text_box(slide5, left, Inches(5.3), Inches(2.8), Inches(0.4),
                     title, font_size=12, font_color=COLORS['gray3'], bold=True)
        add_text_box(slide5, left, Inches(5.7), Inches(2.8), Inches(0.4),
                     sub, font_size=10, font_color=COLORS['gray4'])
    
    # ===================
    # SLIDE 6: Roadmap
    # ===================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6, COLORS['bg'])
    
    add_text_box(slide6, Inches(0.5), Inches(0.5), Inches(3), Inches(0.5),
                 "ROADMAP", font_size=12, font_color=COLORS['lavender'], bold=True)
    
    add_text_box(slide6, Inches(0.5), Inches(1.2), Inches(10), Inches(1),
                 "From testnet to mainnet in 90 days",
                 font_size=36, font_color=COLORS['white'], bold=True)
    
    # Milestones
    milestones = [
        ("30", "Testnet", "APR", 
         "⚡ MPC wallet infrastructure\n⚡ x402 pay-as-you-use\n⚡ Full testnet deployment",
         COLORS['lavender']),
        ("60", "Mainnet", "MAY",
         "✓ Real users, real wallets\n✓ Developer payments live\n✓ Production-grade security",
         COLORS['mint']),
        ("90", "Scale", "JUN",
         "◆ Etherfuse yield in Mexico\n◆ Enterprise migration layer\n◆ EVM↔Stellar bridge scoping",
         COLORS['blush']),
    ]
    
    for i, (days, title, month, items, color) in enumerate(milestones):
        left = Inches(0.5 + i * 4.2)
        # Days circle
        add_text_box(slide6, left, Inches(2.5), Inches(1), Inches(0.6),
                     days, font_size=24, font_color=color, bold=True)
        # Month
        add_text_box(slide6, left, Inches(3.1), Inches(1), Inches(0.4),
                     month, font_size=12, font_color=color, bold=True)
        # Title
        add_text_box(slide6, left, Inches(3.6), Inches(3.5), Inches(0.5),
                     title, font_size=18, font_color=color, bold=True)
        # Items
        add_text_box(slide6, left, Inches(4.2), Inches(3.5), Inches(2),
                     items, font_size=14, font_color=COLORS['white'])
    
    # ===================
    # SLIDE 7: Close
    # ===================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7, COLORS['bg'])
    
    add_text_box(slide7, Inches(0.5), Inches(0.5), Inches(3), Inches(0.5),
                 "ACCESLY", font_size=12, font_color=COLORS['lavender'], bold=True)
    
    add_text_box(slide7, Inches(1), Inches(2.5), Inches(11), Inches(1),
                 "Join us. Build with us.", font_size=48, font_color=COLORS['white'], bold=True)
    
    add_text_box(slide7, Inches(1), Inches(4), Inches(11), Inches(0.5),
                 "Scan to connect", font_size=18, font_color=COLORS['gray4'])
    
    add_text_box(slide7, Inches(5), Inches(5.5), Inches(3), Inches(0.5),
                 "[QR CODE]", font_size=24, font_color=COLORS['gray5'])
    
    add_text_box(slide7, Inches(5.5), Inches(6.5), Inches(2), Inches(0.5),
                 "@accesly", font_size=18, font_color=COLORS['gray3'], bold=True)
    
    # Save
    output_path = os.path.join(os.path.dirname(__file__), 'accesly-pitch.pptx')
    prs.save(output_path)
    print(f"✅ PPTX saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
