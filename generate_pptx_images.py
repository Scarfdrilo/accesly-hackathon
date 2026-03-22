#!/usr/bin/env python3
"""
Generate PPTX from slide images
"""
from pptx import Presentation
from pptx.util import Inches
import os

SLIDES_DIR = os.path.join(os.path.dirname(__file__), 'slides')
OUTPUT_PATH = os.path.join(os.path.dirname(__file__), 'accesly-pitch.pptx')

def create_pptx():
    prs = Presentation()
    # 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # Add each slide image
    for i in range(1, 8):
        img_path = os.path.join(SLIDES_DIR, f'slide_{i}.png')
        if os.path.exists(img_path):
            slide = prs.slides.add_slide(blank_layout)
            # Full-slide image
            slide.shapes.add_picture(
                img_path,
                Inches(0), Inches(0),
                width=prs.slide_width,
                height=prs.slide_height
            )
            print(f"✅ Added slide {i}")
        else:
            print(f"⚠️ Missing: {img_path}")
    
    prs.save(OUTPUT_PATH)
    print(f"\n📁 PPTX saved to: {OUTPUT_PATH}")

if __name__ == "__main__":
    create_pptx()
